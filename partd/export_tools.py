"""
Export Tools for Document Generation
Supports PDF, DOCX, and PPT export
"""

import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent))

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib import colors

import os
from datetime import datetime
from typing import Dict, List, Optional
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches, Pt as PPTPt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPT export will not work.")


class DocumentExporter:
    """Export content to various document formats for Task 2."""
    
    def __init__(self, output_dir: str = "generated_documents"):
        """
        Initialize document exporter.
        
        Args:
            output_dir: Directory to save generated documents
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True, parents=True)
        logger.info(f"✓ Document exporter initialized: {self.output_dir}")
    
    def export_quiz_to_docx(self, quiz_data: Dict) -> str:
        """Export quiz to DOCX."""
        logger.info(f"Exporting quiz to DOCX: {quiz_data['topic']}")
        
        doc = Document()
        
        # Title
        title = doc.add_heading(f"Quiz: {quiz_data['topic'].title()}", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata
        doc.add_paragraph()
        doc.add_paragraph(f"Topic: {quiz_data['topic'].title()}")
        doc.add_paragraph(f"Questions: {quiz_data['num_questions']}")
        doc.add_paragraph(f"Difficulty: {quiz_data['difficulty'].title()}")
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}")
        doc.add_paragraph("_" * 80)
        doc.add_paragraph()
        
        # Questions
        questions = quiz_data.get('questions', [])
        
        for i, q in enumerate(questions, 1):
            doc.add_heading(f"Question {i}", level=2)
            doc.add_paragraph(q['question'])
            doc.add_paragraph()
            
            for key in ['A', 'B', 'C', 'D']:
                if key in q['options']:
                    doc.add_paragraph(f"    {key}) {q['options'][key]}")
            
            doc.add_paragraph()
            doc.add_paragraph()
        
        # Answer Key
        doc.add_page_break()
        doc.add_heading("Answer Key", level=1)
        doc.add_paragraph()
        
        for i, q in enumerate(questions, 1):
            answer_para = doc.add_paragraph()
            answer_para.add_run(f"{i}. ").bold = True
            answer_para.add_run(f"Answer: {q['correct_answer']}")
            
            explanation = doc.add_paragraph(f"   {q['explanation']}")
            explanation.style = 'Intense Quote'
            doc.add_paragraph()
        
        # Save
        filename = f"quiz_{quiz_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        
        logger.info(f"✓ Quiz saved: {filepath}")
        return str(filepath)
    
    def export_quiz_to_pdf(self, quiz_data: Dict) -> str:
        """Export quiz to PDF with proper formatting and encoding."""
        logger.info(f"Exporting quiz to PDF: {quiz_data['topic']}")
        
        # Clean filename - remove special characters
        topic_clean = ''.join(c if c.isalnum() or c in (' ', '-', '_') else '_' for c in quiz_data['topic'])
        filename = f"quiz_{topic_clean.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = self.output_dir / filename
        
        # Helper function to escape HTML and clean text
        def clean_text(text):
            """Clean and escape text for PDF."""
            if not text:
                return ""
            # Replace common problematic characters
            text = str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            # Remove or replace other special characters that might cause issues
            text = text.replace('\x00', '')  # Remove null bytes
            return text
        
        try:
            doc = SimpleDocTemplate(
                str(filepath),
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )
            styles = getSampleStyleSheet()
            story = []
            
            # Title style
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2C5F2D'),
                alignment=1,
                spaceAfter=12
            )
            
            # Add title
            title_text = clean_text(f"Quiz: {quiz_data.get('topic', 'Quiz').title()}")
            story.append(Paragraph(f"<b>{title_text}</b>", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Add metadata
            metadata_style = ParagraphStyle(
                'Metadata',
                parent=styles['Normal'],
                fontSize=10,
                textColor=colors.HexColor('#666666'),
                alignment=1
            )
            story.append(Paragraph(f"Topic: {clean_text(quiz_data.get('topic', 'N/A'))}", metadata_style))
            story.append(Paragraph(f"Questions: {quiz_data.get('num_questions', len(quiz_data.get('questions', [])))}", metadata_style))
            story.append(Paragraph(f"Difficulty: {clean_text(quiz_data.get('difficulty', 'Medium').title())}", metadata_style))
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}", metadata_style))
            story.append(Spacer(1, 0.3*inch))
            
            # Questions
            questions = quiz_data.get('questions', [])
            if not questions:
                story.append(Paragraph("No questions available.", styles['Normal']))
            else:
                for i, q in enumerate(questions, 1):
                    # Question number and text
                    question_text = clean_text(q.get('question', f'Question {i}'))
                    story.append(Paragraph(f"<b>Question {i}</b>", styles['Heading2']))
                    story.append(Paragraph(question_text, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                    
                    # Options
                    options = q.get('options', {})
                    for key in ['A', 'B', 'C', 'D']:
                        if key in options:
                            option_text = clean_text(options[key])
                            story.append(Paragraph(f"&nbsp;&nbsp;&nbsp;&nbsp;<b>{key})</b> {option_text}", styles['Normal']))
                    
                    story.append(Spacer(1, 0.3*inch))
            
            # Answer Key
            story.append(PageBreak())
            answer_title_style = ParagraphStyle(
                'AnswerTitle',
                parent=title_style,
                fontSize=20
            )
            story.append(Paragraph("<b>Answer Key</b>", answer_title_style))
            story.append(Spacer(1, 0.2*inch))
            
            for i, q in enumerate(questions, 1):
                correct_answer = clean_text(q.get('correct_answer', 'N/A'))
                explanation = clean_text(q.get('explanation', 'No explanation provided.'))
                story.append(Paragraph(f"<b>{i}. Answer: {correct_answer}</b>", styles['Normal']))
                story.append(Paragraph(f"<i>Explanation: {explanation}</i>", styles['Italic']))
                story.append(Spacer(1, 0.15*inch))
            
            # Build PDF
            doc.build(story)
            
            # Verify file was created and is readable
            if not filepath.exists():
                raise FileNotFoundError(f"PDF file was not created: {filepath}")
            
            file_size = filepath.stat().st_size
            if file_size == 0:
                raise ValueError(f"PDF file is empty: {filepath}")
            
            logger.info(f"✓ Quiz PDF saved: {filepath} ({file_size} bytes)")
            return str(filepath)
            
        except Exception as e:
            logger.error(f"Failed to export quiz to PDF: {e}")
            import traceback
            logger.error(traceback.format_exc())
            # Try to clean up corrupted file
            if filepath.exists():
                try:
                    filepath.unlink()
                except:
                    pass
            raise
    
    def export_study_guide_to_docx(self, guide_data: Dict) -> str:
        """Export study guide to DOCX."""
        logger.info(f"Exporting study guide to DOCX: {guide_data['topic']}")
        
        doc = Document()
        
        # Title
        title = doc.add_heading(f"Study Guide: {guide_data['topic'].title()}", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        subtitle = doc.add_paragraph(f"Environmental Sustainability Education")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.style.font.italic = True
        
        doc.add_paragraph()
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_paragraph("_" * 80)
        doc.add_paragraph()
        
        # Content
        content = guide_data.get('content', '')
        
        for line in content.split('\n'):
            line = line.strip()
            
            if not line:
                doc.add_paragraph()
                continue
            
            if line.startswith('##'):
                doc.add_heading(line.replace('##', '').strip(), level=2)
            elif line.startswith('#'):
                doc.add_heading(line.replace('#', '').strip(), level=1)
            elif line.startswith('-') or line.startswith('•'):
                doc.add_paragraph(line[1:].strip(), style='List Bullet')
            elif line[0].isdigit() and '. ' in line[:5]:
                doc.add_paragraph(line, style='List Number')
            else:
                doc.add_paragraph(line)
        
        # Save
        filename = f"study_guide_{guide_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        
        logger.info(f"✓ Study guide saved: {filepath}")
        return str(filepath)
    
    def export_study_guide_to_pdf(self, guide_data: Dict) -> str:
        """Export study guide to PDF."""
        logger.info(f"Exporting study guide to PDF: {guide_data['topic']}")
        
        filename = f"study_guide_{guide_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = self.output_dir / filename
        
        doc = SimpleDocTemplate(str(filepath), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#2c5f2d'),
            spaceAfter=30,
            alignment=1  # Center
        )
        story.append(Paragraph(f"Study Guide: {guide_data['topic'].title()}", title_style))
        story.append(Spacer(1, 0.2*inch))
        
        subtitle_style = ParagraphStyle(
            'CustomSubtitle',
            parent=styles['Normal'],
            fontSize=12,
            fontName='Helvetica-Oblique',
            alignment=1  # Center
        )
        story.append(Paragraph("Environmental Sustainability Education", subtitle_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Metadata
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph("_" * 80, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Content
        content = guide_data.get('content', '')
        
        for line in content.split('\n'):
            line = line.strip()
            
            if not line:
                story.append(Spacer(1, 0.1*inch))
                continue
            
            if line.startswith('##'):
                heading_text = line.replace('##', '').strip()
                story.append(Paragraph(heading_text, styles['Heading2']))
                story.append(Spacer(1, 0.15*inch))
            elif line.startswith('#'):
                heading_text = line.replace('#', '').strip()
                story.append(Paragraph(heading_text, styles['Heading1']))
                story.append(Spacer(1, 0.2*inch))
            elif line.startswith('-') or line.startswith('•'):
                bullet_text = line[1:].strip()
                story.append(Paragraph(f"• {bullet_text}", styles['Normal']))
                story.append(Spacer(1, 0.05*inch))
            elif line[0].isdigit() and '. ' in line[:5]:
                story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 0.05*inch))
            else:
                story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
        
        doc.build(story)
        logger.info(f"✓ Study guide PDF saved: {filepath}")
        return str(filepath)
    
    def export_study_guide_to_ppt(self, guide_data: Dict) -> str:
        """Export study guide to PPT."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Cannot export to PPT.")
        
        logger.info(f"Exporting study guide to PPT: {guide_data['topic']}")
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Study Guide: {guide_data['topic'].title()}"
        subtitle.text = "Environmental Sustainability Education"
        
        # Content slides
        content = guide_data.get('content', '')
        current_slide = None
        current_content = []
        
        for line in content.split('\n'):
            line = line.strip()
            
            if not line:
                continue
            
            if line.startswith('##') or line.startswith('#'):
                # Save previous slide if exists
                if current_slide and current_content:
                    content_text = '\n'.join(current_content)
                    if len(content_text) > 0:
                        tf = current_slide.placeholders[1].text_frame
                        tf.text = content_text
                
                # New slide for heading
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                title = current_slide.shapes.title
                title.text = line.replace('#', '').strip()
                current_content = []
            else:
                current_content.append(line)
        
        # Save last slide
        if current_slide and current_content:
            content_text = '\n'.join(current_content)
            if len(content_text) > 0:
                tf = current_slide.placeholders[1].text_frame
                tf.text = content_text
        
        # Save
        filename = f"study_guide_{guide_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        logger.info(f"✓ Study guide PPT saved: {filepath}")
        return str(filepath)
    
    def export_content_to_docx(self, content_data: Dict, content_type: str = "content") -> str:
        """Export general content to DOCX."""
        topic = content_data.get('topic', 'Content')
        
        doc = Document()
        doc.add_heading(topic.title(), 0)
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_paragraph("_" * 80)
        doc.add_paragraph()
        
        content = content_data.get('content', '')
        
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                doc.add_paragraph()
                continue
            
            if line.startswith('##'):
                doc.add_heading(line.replace('##', '').strip(), level=2)
            elif line.startswith('#'):
                doc.add_heading(line.replace('#', '').strip(), level=1)
            elif line.startswith('-') or line.startswith('•'):
                doc.add_paragraph(line[1:].strip(), style='List Bullet')
            else:
                doc.add_paragraph(line)
        
        topic_safe = content_data.get('topic', 'content').replace(' ', '_')
        filename = f"{content_type}_{topic_safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        
        logger.info(f"✓ Content saved: {filepath}")
        return str(filepath)
    
    def export_quiz_to_ppt(self, quiz_data: Dict) -> str:
        """Export quiz to PPT (PowerPoint)."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        logger.info(f"Exporting quiz to PPT: {quiz_data['topic']}")
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Quiz: {quiz_data['topic'].title()}"
        subtitle.text = f"{quiz_data['num_questions']} Questions | Difficulty: {quiz_data['difficulty'].title()}"
        
        # Questions
        questions = quiz_data.get('questions', [])
        
        for i, q in enumerate(questions, 1):
            # Question slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"Question {i}"
            tf = body_shape.text_frame
            tf.text = q['question']
            
            # Add options
            for key in ['A', 'B', 'C', 'D']:
                if key in q['options']:
                    p = tf.add_paragraph()
                    p.text = f"{key}) {q['options'][key]}"
                    p.level = 1
        
        # Answer key slide
        answer_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(answer_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Answer Key"
        tf = body_shape.text_frame
        
        for i, q in enumerate(questions, 1):
            p = tf.add_paragraph()
            p.text = f"Question {i}: {q['correct_answer']}"
            p.level = 0
            
            p2 = tf.add_paragraph()
            p2.text = q['explanation']
            p2.level = 1
        
        # Save
        filename = f"quiz_{quiz_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        logger.info(f"✓ Quiz PPT saved: {filepath}")
        return str(filepath)
    
    def export_content_to_ppt(self, content_data: Dict, content_type: str = "content") -> str:
        """Export general content to PPT."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        topic = content_data.get('topic', 'Content')
        logger.info(f"Exporting {content_type} to PPT: {topic}")
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic.title()
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Content slides
        content = content_data.get('content', '')
        paragraphs = content.split('\n')
        
        current_slide = None
        current_tf = None
        lines_on_slide = 0
        max_lines_per_slide = 8
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # New slide for headings
            if para.startswith('#') or lines_on_slide >= max_lines_per_slide:
                if current_slide is None or lines_on_slide >= max_lines_per_slide:
                    bullet_slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = current_slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    
                    if para.startswith('#'):
                        title_shape.text = para.replace('#', '').strip()
                        current_tf = body_shape.text_frame
                        current_tf.text = ""
                    else:
                        title_shape.text = topic.title()
                        current_tf = body_shape.text_frame
                        current_tf.text = para
                    lines_on_slide = 1
                    continue
            
            if current_slide is None:
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = topic.title()
                current_tf = body_shape.text_frame
                current_tf.text = ""
            
            p = current_tf.add_paragraph()
            if para.startswith('-') or para.startswith('•'):
                p.text = para[1:].strip()
                p.level = 1
            else:
                p.text = para
                p.level = 0
            lines_on_slide += 1
        
        # Save
        topic_safe = topic.replace(' ', '_')
        filename = f"{content_type}_{topic_safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        logger.info(f"✓ Content PPT saved: {filepath}")
        return str(filepath)

        topic_safe = content_data.get('topic', 'content').replace(' ', '_')
        filename = f"{content_type}_{topic_safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = self.output_dir / filename
        doc.save(str(filepath))
        
        logger.info(f"✓ Content saved: {filepath}")
        return str(filepath)
    
    def export_quiz_to_ppt(self, quiz_data: Dict) -> str:
        """Export quiz to PPT (PowerPoint)."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        logger.info(f"Exporting quiz to PPT: {quiz_data['topic']}")
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Quiz: {quiz_data['topic'].title()}"
        subtitle.text = f"{quiz_data['num_questions']} Questions | Difficulty: {quiz_data['difficulty'].title()}"
        
        # Questions
        questions = quiz_data.get('questions', [])
        
        for i, q in enumerate(questions, 1):
            # Question slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"Question {i}"
            tf = body_shape.text_frame
            tf.text = q['question']
            
            # Add options
            for key in ['A', 'B', 'C', 'D']:
                if key in q['options']:
                    p = tf.add_paragraph()
                    p.text = f"{key}) {q['options'][key]}"
                    p.level = 1
        
        # Answer key slide
        answer_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(answer_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Answer Key"
        tf = body_shape.text_frame
        
        for i, q in enumerate(questions, 1):
            p = tf.add_paragraph()
            p.text = f"Question {i}: {q['correct_answer']}"
            p.level = 0
            
            p2 = tf.add_paragraph()
            p2.text = q['explanation']
            p2.level = 1
        
        # Save
        filename = f"quiz_{quiz_data['topic'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        logger.info(f"✓ Quiz PPT saved: {filepath}")
        return str(filepath)
    
    def export_content_to_ppt(self, content_data: Dict, content_type: str = "content") -> str:
        """Export general content to PPT."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        topic = content_data.get('topic', 'Content')
        logger.info(f"Exporting {content_type} to PPT: {topic}")
        
        prs = Presentation()
        prs.slide_width = PPTInches(10)
        prs.slide_height = PPTInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic.title()
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Content slides
        content = content_data.get('content', '')
        paragraphs = content.split('\n')
        
        current_slide = None
        current_tf = None
        lines_on_slide = 0
        max_lines_per_slide = 8
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # New slide for headings
            if para.startswith('#') or lines_on_slide >= max_lines_per_slide:
                if current_slide is None or lines_on_slide >= max_lines_per_slide:
                    bullet_slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = current_slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    
                    if para.startswith('#'):
                        title_shape.text = para.replace('#', '').strip()
                        current_tf = body_shape.text_frame
                        current_tf.text = ""
                    else:
                        title_shape.text = topic.title()
                        current_tf = body_shape.text_frame
                        current_tf.text = para
                    lines_on_slide = 1
                    continue
            
            if current_slide is None:
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = topic.title()
                current_tf = body_shape.text_frame
                current_tf.text = ""
            
            p = current_tf.add_paragraph()
            if para.startswith('-') or para.startswith('•'):
                p.text = para[1:].strip()
                p.level = 1
            else:
                p.text = para
                p.level = 0
            lines_on_slide += 1
        
        # Save
        topic_safe = topic.replace(' ', '_')
        filename = f"{content_type}_{topic_safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = self.output_dir / filename
        prs.save(str(filepath))
        
        logger.info(f"✓ Content PPT saved: {filepath}")
        return str(filepath)
